#!/usr/bin/env python3
"""
Auto-generate PowerPoint dashboard from Excel files
"""

import openpyxl
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime, timedelta

def extract_health_data(filename):
    """Extract data from Health Tracker Excel"""
    wb = openpyxl.load_workbook(filename, data_only=True)
    feb_sheet = wb['Feb']
    
    # Row 34 = Totals
    total_loaded = feb_sheet.cell(34, 8).value
    total_loaded_vs_drop = feb_sheet.cell(34, 13).value
    total_late_loads = feb_sheet.cell(34, 16).value
    total_roll_over = feb_sheet.cell(34, 17).value
    
    # Count days
    day_count = 0
    for row_idx in range(6, 34):
        val = feb_sheet.cell(row_idx, 8).value
        cell_b = feb_sheet.cell(row_idx, 2).value
        if cell_b and 'total' in str(cell_b).lower():
            continue
        if val and val != 0:
            day_count += 1
    
    # Bin statistics
    bin_values = []
    for row_idx in range(6, 34):
        cell_b = feb_sheet.cell(row_idx, 2).value
        if cell_b and 'total' in str(cell_b).lower():
            continue
        bins = feb_sheet.cell(row_idx, 19).value
        if bins and isinstance(bins, (int, float)) and bins > 0:
            bin_values.append(int(bins))
    
    return {
        'running_total': int(total_loaded) if total_loaded else 0,
        'days_tracked': day_count,
        'avg_per_day': int(total_loaded / day_count) if day_count > 0 and total_loaded else 0,
        'loaded_vs_drop_total': int(total_loaded_vs_drop) if total_loaded_vs_drop else 0,
        'avg_loaded_vs_drop': int(total_loaded_vs_drop / day_count) if day_count > 0 and total_loaded_vs_drop else 0,
        'late_loads': int(total_late_loads) if total_late_loads else 0,
        'avg_bins': int(sum(bin_values) / len(bin_values)) if bin_values else 0,
        'max_bins': max(bin_values) if bin_values else 0,
        'min_bins': min(bin_values) if bin_values else 0,
    }

def extract_picker_data(filename):
    """Extract picker efficiency data"""
    try:
        wb = openpyxl.load_workbook(filename)
        sheet = wb['Formulas']
        total_picked = 0
        total_hours = 0
        count = 0
        
        for row in sheet.iter_rows(min_row=3, max_row=30, values_only=True):
            if row[0] and row[1]:
                hours = row[2] if row[2] else 0
                pallets = sum([row[i] for i in range(3, 10) if row[i] and isinstance(row[i], (int, float))])
                if pallets > 0 or hours > 0:
                    total_picked += pallets
                    total_hours += hours
                    count += 1
        
        efficiency = round(total_picked / total_hours, 2) if total_hours > 0 else 0
        return {'total': int(total_picked), 'count': count, 'efficiency': efficiency}
    except:
        return {'total': 0, 'count': 0, 'efficiency': 0}

def extract_putaway_data(filename):
    """Extract putaway efficiency data"""
    try:
        wb = openpyxl.load_workbook(filename)
        sheet = wb['Formulas']
        total_putaway = 0
        total_hours = 0
        count = 0
        
        for row in sheet.iter_rows(min_row=3, max_row=30, values_only=True):
            if row[0] and row[1]:
                hours = row[2] if row[2] else 0
                pallets = sum([row[i] for i in range(3, 10) if row[i] and isinstance(row[i], (int, float))])
                if pallets > 0 or hours > 0:
                    total_putaway += pallets
                    total_hours += hours
                    count += 1
        
        efficiency = round(total_putaway / total_hours, 2) if total_hours > 0 else 0
        return {'total': int(total_putaway), 'count': count, 'efficiency': efficiency}
    except:
        return {'total': 0, 'count': 0, 'efficiency': 0}

def calculate_safety_days():
    """Calculate days since Nov 29, 2023"""
    start_date = datetime(2023, 11, 29)
    today = datetime.now()
    return (today - start_date).days

def create_presentation(health, picker, putaway, safety_days):
    """Create PowerPoint presentation"""
    
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # Colors
    bg_color = RGBColor(26, 26, 46)
    accent_color = RGBColor(78, 205, 196)
    green_color = RGBColor(11, 232, 129)
    white_color = RGBColor(255, 255, 255)
    
    def add_title_slide(title, subtitle):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
        
        title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(2))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(72)
        title_para.font.bold = True
        title_para.font.color.rgb = accent_color
        title_para.alignment = PP_ALIGN.CENTER
        
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(14), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(36)
        subtitle_para.font.color.rgb = white_color
        subtitle_para.alignment = PP_ALIGN.CENTER
    
    def add_kpi_slide(title, kpis):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
        
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.color.rgb = accent_color
        title_para.alignment = PP_ALIGN.CENTER
        
        num_kpis = len(kpis)
        cols = 2
        card_width = Inches(6.5)
        card_height = Inches(2.5)
        start_x = Inches(1.5)
        start_y = Inches(2)
        gap = Inches(0.5)
        
        for idx, kpi in enumerate(kpis):
            row = idx // cols
            col = idx % cols
            x = start_x + col * (card_width + gap)
            y = start_y + row * (card_height + gap)
            
            card = slide.shapes.add_shape(1, x, y, card_width, card_height)
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(40, 40, 70)
            card.line.color.rgb = accent_color
            card.line.width = Pt(3)
            
            label_box = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.3), card_width - Inches(0.6), Inches(0.7))
            label_frame = label_box.text_frame
            label_frame.text = kpi['label']
            label_para = label_frame.paragraphs[0]
            label_para.font.size = Pt(24)
            label_para.font.color.rgb = RGBColor(160, 160, 160)
            label_para.alignment = PP_ALIGN.CENTER
            
            value_box = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(1), card_width - Inches(0.6), Inches(1))
            value_frame = value_box.text_frame
            value_frame.text = kpi['value']
            value_para = value_frame.paragraphs[0]
            value_para.font.size = Pt(56)
            value_para.font.bold = True
            value_para.font.color.rgb = green_color if kpi.get('highlight') else white_color
            value_para.alignment = PP_ALIGN.CENTER
            
            if 'sublabel' in kpi:
                sub_box = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(2), card_width - Inches(0.6), Inches(0.4))
                sub_frame = sub_box.text_frame
                sub_frame.text = kpi['sublabel']
                sub_para = sub_frame.paragraphs[0]
                sub_para.font.size = Pt(18)
                sub_para.font.color.rgb = RGBColor(120, 120, 120)
                sub_para.alignment = PP_ALIGN.CENTER
    
    # Create slides
    add_title_slide("US5 Warehouse Dashboard", f"Updated: {datetime.now().strftime('%B %d, %Y')}")
    
    add_kpi_slide("Safety & Loader Operations", [
        {'label': 'Days Without OSHA Recordable', 'value': str(safety_days), 'highlight': True},
        {'label': 'Total Tires Loaded', 'value': f'{health["running_total"]:,}', 'sublabel': f'{health["days_tracked"]} days tracked'},
        {'label': 'Average Per Day', 'value': f'{health["avg_per_day"]:,}'},
        {'label': 'Loaded vs Drop Avg', 'value': f'{health["avg_loaded_vs_drop"]:,}/day'},
    ])
    
    add_kpi_slide("Warehouse Operations", [
        {'label': 'Pallets Picked', 'value': f'{picker["total"]:,}', 'sublabel': f'{picker["count"]} pickers'},
        {'label': 'Pallets Putaway', 'value': f'{putaway["total"]:,}', 'sublabel': f'{putaway["count"]} staff'},
        {'label': 'Picker Efficiency', 'value': f'{picker["efficiency"]}', 'sublabel': 'pallets per hour'},
        {'label': 'Putaway Efficiency', 'value': f'{putaway["efficiency"]}', 'sublabel': 'pallets per hour'},
    ])
    
    add_kpi_slide("Additional Metrics", [
        {'label': 'Late Loads', 'value': str(health["late_loads"]), 'highlight': health["late_loads"] == 0},
        {'label': 'Average Bins Available', 'value': str(health["avg_bins"])},
        {'label': 'Max Bins', 'value': str(health["max_bins"])},
        {'label': 'Min Bins', 'value': str(health["min_bins"])},
    ])
    
    add_kpi_slide("Employee of the Month - February 2026", [
        {'label': 'Loader of the Month', 'value': 'Malachi Burgess', 'highlight': True},
        {'label': 'Picker of the Month', 'value': 'Makaela Culler', 'highlight': True},
        {'label': 'Putaway of the Month', 'value': 'Donald Lynch', 'highlight': True},
    ])
    
    return prs

if __name__ == '__main__':
    print("Extracting data from Excel files...")
    
    health = extract_health_data('Health_Tracker_2026_xlsx.xlsx')
    picker = extract_picker_data('Picker_Efficiency_2026_xlsx.xlsx')
    putaway = extract_putaway_data('Putaway_Efficiency_2026_xlsx.xlsx')
    safety_days = calculate_safety_days()
    
    print(f"✓ Loader: {health['running_total']:,} tires ({health['days_tracked']} days)")
    print(f"✓ Picker: {picker['total']:,} pallets")
    print(f"✓ Putaway: {putaway['total']:,} pallets")
    print(f"✓ Safety: {safety_days} days")
    
    print("\nGenerating PowerPoint presentation...")
    prs = create_presentation(health, picker, putaway, safety_days)
    prs.save('US5_Warehouse_KPI_Dashboard.pptx')
    
    print("✓ PowerPoint created: US5_Warehouse_KPI_Dashboard.pptx")
